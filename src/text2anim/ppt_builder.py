"""PPT 合成模块：将图片合成为 PowerPoint 幻灯片"""

import logging
from pathlib import Path

import questionary
from pptx import Presentation
from pptx.util import Inches, Pt

from text2anim.config import IMAGES_DIR, LOGS_DIR, PPT_DIR

logger = logging.getLogger(__name__)

# 支持的图片格式
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}


def setup_logging() -> None:
    """配置日志输出到 logs/ 目录"""
    LOGS_DIR.mkdir(parents=True, exist_ok=True)
    log_file = LOGS_DIR / "generate_ppt.log"
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        handlers=[
            logging.FileHandler(log_file, encoding="utf-8"),
            logging.StreamHandler(),
        ],
    )


def list_image_folders() -> list[Path]:
    """列出 images/ 下所有包含图片的子文件夹"""
    if not IMAGES_DIR.exists():
        return []

    folders = []
    for item in sorted(IMAGES_DIR.iterdir()):
        if not item.is_dir():
            continue
        # 检查是否包含图片文件
        has_images = any(
            f.suffix.lower() in IMAGE_EXTENSIONS for f in item.iterdir()
        )
        if has_images:
            folders.append(item)

    return folders


def select_folder(folder_name: str | None) -> Path:
    """选择要使用的图片子文件夹"""
    folders = list_image_folders()

    if not folders:
        raise FileNotFoundError(
            f"images/ 目录下没有找到包含图片的子文件夹\n"
            f"请先运行 generate-images 生成图片"
        )

    # 如果指定了文件夹名称，直接查找
    if folder_name:
        for f in folders:
            if f.name == folder_name:
                return f
        raise FileNotFoundError(
            f"找不到子文件夹: {folder_name}\n"
            f"可用的子文件夹: {[f.name for f in folders]}"
        )

    # 只有一个文件夹时直接使用
    if len(folders) == 1:
        logger.info("自动选择: %s", folders[0].name)
        return folders[0]

    # 多个文件夹时交互式选择（questionary 光标选择）
    choices = []
    for f in folders:
        image_count = sum(
            1 for file in f.iterdir()
            if file.suffix.lower() in IMAGE_EXTENSIONS
        )
        choices.append(questionary.Choice(f"{f.name} ({image_count} 张图片)", value=f))

    selected = questionary.select("选择图片文件夹:", choices=choices).ask()
    if selected is None:
        raise KeyboardInterrupt("用户取消了选择")
    return selected


def get_images(folder: Path) -> list[Path]:
    """获取文件夹中所有图片文件，按文件名排序"""
    images = [
        f for f in folder.iterdir()
        if f.suffix.lower() in IMAGE_EXTENSIONS
    ]
    return sorted(images, key=lambda p: p.name)


def extract_title(filename: str) -> str:
    """从文件名中提取标题（去掉序号前缀和扩展名）"""
    name = Path(filename).stem
    # 去掉 "01_" 这样的序号前缀
    if len(name) > 2 and name[:2].isdigit() and name[2] == "_":
        name = name[3:]
    return name


def build_ppt(images: list[Path], output_path: Path) -> None:
    """创建 PPT 文件"""
    prs = Presentation()
    # 使用 16:9 宽屏布局
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 使用空白布局
    blank_layout = prs.slide_layouts[6]

    for image_path in images:
        slide = prs.slides.add_slide(blank_layout)

        # 图片居中放置，占幻灯片大部分区域
        # 计算图片尺寸（保持宽高比，最大 10x6 英寸）
        img_width = Inches(10)
        img_left = Inches(1.667)  # (13.333 - 10) / 2
        img_top = Inches(0.5)

        pic = slide.shapes.add_picture(
            str(image_path), img_left, img_top, width=img_width,
        )

        # 如果图片太高，按比例缩小
        if pic.height > Inches(6):
            ratio = Inches(6) / pic.height
            pic.height = Inches(6)
            pic.width = int(pic.width * ratio)
            # 重新居中
            pic.left = int((prs.slide_width - pic.width) / 2)

        # 底部添加标题文字
        title = extract_title(image_path.name)
        if title:
            left = Inches(1)
            top = Inches(6.8)
            width = Inches(11.333)
            height = Inches(0.5)
            txbox = slide.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.color.rgb = None  # 使用默认颜色

    prs.save(str(output_path))


def run_generate_ppt(folder_name: str | None) -> None:
    """CLI 入口：生成 PPT"""
    setup_logging()

    folder = select_folder(folder_name)
    images = get_images(folder)

    if not images:
        raise ValueError(f"子文件夹中没有找到图片: {folder}")

    logger.info("选中文件夹: %s (%d 张图片)", folder.name, len(images))

    # 输出 PPT
    PPT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = PPT_DIR / f"{folder.name}.pptx"

    build_ppt(images, output_path)
    logger.info("PPT 已生成: %s", output_path)
